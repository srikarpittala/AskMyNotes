import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from transformers import AutoTokenizer, AutoModelForCausalLM, pipeline

# === UI Header ===
st.title("📄 AskMyNotes")
st.write("Upload a document and ask your question below.")

# File Upload
uploaded_file = st.file_uploader("Upload file (PDF, DOCX, PPTX, TXT)", type=["pdf", "txt", "docx", "pptx"])

# Text Extraction
def extract_text(file):
    file_type = file.name.split('.')[-1].lower()
    if file_type == 'pdf':
        reader = PdfReader(file)
        return " ".join(page.extract_text() or "" for page in reader.pages)
    elif file_type == 'txt':
        return file.read().decode("utf-8")
    elif file_type == 'docx':
        doc = Document(file)
        return " ".join(para.text for para in doc.paragraphs)
    elif file_type == 'pptx':
        prs = Presentation(file)
        return " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    return ""

# Main Logic
if uploaded_file is not None:
    text = extract_text(uploaded_file)

    if text.strip():
        # Split & Embed
        splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)
        chunks = splitter.split_text(text)

        @st.cache_resource
        def get_store(data_chunks):
            embed = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
            return FAISS.from_texts(data_chunks, embed)

        store = get_store(chunks)

        @st.cache_resource
        def load_model():
            model_name = "TinyLlama/TinyLlama-1.1B-Chat-v1.0"
            tokenizer = AutoTokenizer.from_pretrained(model_name)
            if tokenizer.pad_token is None:
                tokenizer.pad_token = tokenizer.eos_token
            model = AutoModelForCausalLM.from_pretrained(model_name)
            return pipeline("text-generation", model=model, tokenizer=tokenizer)

        gen_pipeline = load_model()

        # Input Box
        question = st.text_input("❓ Ask a question about your document:")

        if question:
            SIMILARITY_THRESHOLD = 1.4 
            docs_with_scores = store.similarity_search_with_score(question, k=7)
            relevant_docs = [doc for doc, score in docs_with_scores if score <= SIMILARITY_THRESHOLD]

            if relevant_docs:
                context = " ".join(doc.page_content for doc in relevant_docs)
                prompt = f"""
You are a helpful and precise assistant. Answer the question using only the content from the document context below.

Strictly adhere to the following rules:
- Do not include outside knowledge.
- Do not guess, invent examples, or hallucinate information.
- Only rephrase or summarize what is explicitly and clearly stated in the document.
- If the document does not contain sufficient information or a definition for the term, respond with: "The document does not provide a definition for this."
- Provide a complete, concise, and direct answer.
- Do not generate any additional questions, conversational filler, or unrelated information after your answer.
- End your answer clearly and stop.

Document Context:
{context}

Question: {question}

Answer:"""

                
                
                output = gen_pipeline(
                    prompt,
                    max_new_tokens=200, 
                                        
                    do_sample=False,
                )[0]["generated_text"]

                # Post-processing the generated text
                
                answer_only = output.replace(prompt, "").strip()

                
                
                unwanted_patterns = ["\nQuestion:", "Document Context:", "Answer:", "Question:"]
                for pattern in unwanted_patterns:
                    if pattern in answer_only:
                        
                        answer_only = answer_only.split(pattern, 1)[0].strip()

                
                lines = answer_only.split('\n')
                unique_lines = []
                for line in lines:
                    stripped_line = line.strip()
                    if stripped_line and stripped_line not in unique_lines:
                        unique_lines.append(stripped_line)

                final_answer = "\n".join(unique_lines)

                
                if final_answer.lower().startswith("answer:"):
                    final_answer = final_answer[len("answer:"):].strip()

                st.subheader("📝 Answer:")
                st.write(final_answer)
            else:
                st.subheader("📝 Answer:")
                st.write("⚠ The document does not contain enough information to answer this question.")
    else:
        st.error("⚠ Could not extract text. Please try another file.")

else:
    st.info("👆 Upload a file to begin.")

    




















    
    








